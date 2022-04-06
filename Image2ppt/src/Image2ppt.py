import os, os.path
from pptx import Presentation
from pptx.util import Inches, Pt
import tkinter
from tkinter import filedialog
from tkinter import *
from tkinter import ttk
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import io
from PIL import Image
from math import ceil
import math
import time

# Utilizes python-pptx: https://python-pptx.readthedocs.io/


def main():

    create_gui = CreateGUI()
    create_gui.construct_form()


class Components:

    def create_textbox(self,frame,text,row,column):
        txt = ttk.Entry(frame,width=40)
        txt.insert(tkinter.END,text)
        txt.grid(row=row,column=column)
        return txt

    def create_label(self,frame,text,row,column):
        label1 = ttk.Label(
            frame,
            text=text,
            padding=(5, 10))
        label1.grid(row=row, column=column)
        return label1

    def create_button(self,frame, text, row, column):
        button1 = ttk.Button(
            frame,
            text=text,
            )
        button1.grid(row=row, column=column)
        return button1

    def create_frame(self,root):
        frame = ttk.Frame(root, padding=40)
        frame.grid()
        return frame

class CreateGUI:
    def __init__(self):

        self.input_path_label = "initial path"
        self.output_path_label = "initial path"
        self.ppt_name_textbox = "test"
        self.path_list = [r"C:\Users\Fridge\Documents\PYGit\Image2ppt\Image2ppt\Input",
                          r"C:\Users\Fridge\Documents\PYGit\Image2ppt\Image2ppt\Output"]

    def construct_form(self):
        components = Components()
        root = Tk()
        root.minsize(width=500, height=300)
        root.title("Image2ppt")
        frame = components.create_frame(root)

        self.input_path_label = components.create_label(frame, self.path_list[0], 0, 0)

        input_path_button = components.create_button(frame, "Input path", 0, 1)
        input_path_button.bind("<ButtonPress>", lambda event: self.get_path(event, self.input_path_label))

        self.output_path_label = components.create_label(frame, self.path_list[1], 1, 0)

        output_path_button = components.create_button(frame,"Output path", 1, 1)
        output_path_button.bind("<ButtonPress>", lambda event: self.get_path(event, self.output_path_label))

        gui_column_desc = components.create_label(frame, "Column Number:", 2, 0)
        self.gui_column = components.create_textbox(frame, 4, 2, 1)

        gui_row_desc = components.create_label(frame, "Row Number:", 3, 0)
        self.gui_row = components.create_textbox(frame, 2, 3, 1)

        gui_ppt_width_desc = components.create_label(frame, "Slide Width:", 4, 0)
        self.gui_ppt_width = components.create_textbox(frame, 13.333, 4, 1)

        gui_ppt_height_desc = components.create_label(frame, "Slide Height:", 5, 0)
        self.gui_ppt_height = components.create_textbox(frame, 7.5, 5, 1)

        gui_slide_counter_desc = components.create_label(frame, "Images for each cell:", 6, 0)
        self.gui_slide_counter = components.create_textbox(frame, 16, 6, 1)

        ppt_name_label = components.create_label (frame, "Save Name:", 7, 0)
        self.gui_ppt_name_textbox = components.create_textbox(frame, "test", 7, 1)

        start_process_button = components.create_button(frame, "Start", 8, 1)
        start_process_button.bind("<ButtonPress>", lambda event: self.ppt_generation_process(event))

        self.start_gui(root)


    def get_path(self, event, arg):
        selected_path = filedialog.askdirectory()
        if not selected_path=="":
            arg.configure(text=selected_path)

    def ppt_generation_process(self, event):
        #tkinter.Tk().withdraw()
        path_list = [self.input_path_label.cget("text"), self.output_path_label.cget("text")]
        input_path = path_list[0]
        output_path = path_list[1]
        parameters = [self.gui_column.get(), self.gui_row.get(), self.gui_slide_counter.get(), self.gui_ppt_width.get(), self.gui_ppt_height.get()]

        append_slide = AppendSlide(input_path)
        append_slide.get_parameters(parameters)
        prs = append_slide.append_images_in_ppt()
        prs.save(output_path + '/' + self.gui_ppt_name_textbox.get() + '.pptx')
        os.startfile(output_path + '/' + self.gui_ppt_name_textbox.get() + '.pptx')

    def start_gui(self,root):
        root.mainloop()


#this is not used in production
class AddTest:
    def add_test(self,a,b):
        return a+b


class AppendSlide:
    # initial value loading
    def __init__(self, input_path):
        self.input_path = input_path

    def get_parameters(self, parameters):
        self.img_list = self.get_images(self.input_path)
        self.img_count = len(self.img_list)
        self.slide_number = 1

        self.column = int(parameters[0])
        self.row =  int(parameters[1])
        self.ppt_width =  float(parameters[3])
        self.ppt_height = float(parameters[4])
        self.img_width = self.ppt_width / self.column
        self.img_height = self.ppt_height / self.row
        self.img_iter = self.column * self.row
        self.slide_counter = int(parameters[2]) / self.img_iter

    def get_images(self, input_path):
        #folder_files = os.listdir(input_path)
        img_list = []
        dir_name = input_path
        # Get list of all files only in the given directory
        list_of_files = filter(lambda x: os.path.isfile(os.path.join(dir_name, x)),
                               os.listdir(dir_name))
        # Sort list of files based on last modification time in ascending order
        list_of_files = sorted(list_of_files,
                               key=lambda x: os.path.getmtime(os.path.join(dir_name, x))
                               , reverse=True
                               )
        # Iterate over sorted list of files and print file path
        # along with last modification time of file
        for file_name in list_of_files:
            file = os.path.join(dir_name, file_name)
            file = file_name
            print(file)
            if file.endswith('.png') or file.endswith(".tif") or file.endswith(".jpg") or file.endswith(".jpeg"):
                 img_list.append(file)

        print(img_list)
        return img_list

    # generate ppt and add images to the ppt
    def append_images_in_ppt(self):
        prs = Presentation()
        prs.slide_width = Inches(self.ppt_width)
        prs.slide_height = Inches(self.ppt_height)
        prs = self.append_images(prs)
        return prs

    def append_images(self, prs):

        self.img_list = self.get_images(self.input_path)
        self.img_count = len(self.img_list)

        blank_slide = prs.slide_layouts[6]
        # in pixels
        pixel_width = int(960 / self.column)
        pixel_height = int(540 / self.row)
        # in inches
        width = self.ppt_width / self.column
        height = self.ppt_height / self.row
        dpi = 72

        for i in range(self.img_count):
            if i % self.img_iter == 0:
                image_slide = prs.slides.add_slide(blank_slide)
                self.textbox(image_slide)

            current_img = Image.open(self.input_path + '/' + self.img_list[i])
            # working in pixels
            ratio = self.get_resize_ratio(current_img.width, current_img.height, pixel_width, pixel_height)
            resized_img = current_img.resize((int(current_img.width * ratio), int(current_img.height * ratio)))

            margin_width = self.get_margin(width,resized_img.width,dpi)
            margin_height = self.get_margin(height,resized_img.height,dpi)
            # in inches
            horizontal = Inches((i % self.column) * (self.ppt_width / self.column))
            vertical = Inches((i % self.img_iter // self.column) * self.ppt_height / self.row)

            vertical_position = self.apply_vertical_margin(i,self.row,self.column,vertical,margin_height)
            horizontal_position = horizontal + margin_width

            with io.BytesIO() as output:
                resized_img.save(output, format="GIF")
                image_slide.shapes.add_picture(output, horizontal_position, vertical_position)

        self.draw_rectangle(image_slide)

        return prs
    def get_margin(self,length,resized_length,dpi):
        return Inches(length - resized_length / dpi) / 2
    def textbox(self,image_slide):
        cellnumber = str(ceil(self.slide_number / self.slide_counter))
        central_box = image_slide.shapes.add_textbox(Inches(self.ppt_width / 2 - 0.65),
                                                     Inches(self.ppt_height / 2 - 0.6), Inches(1), Inches(1))
        central_label = central_box.text_frame.add_paragraph()
        central_label.text = "Cell " + cellnumber
        central_label.font.size = Pt(30)
        self.slide_number += 1
        print(self.slide_number)

    def apply_vertical_margin(self,index,row,column,vertical,margin_height):
        # apply margin
        iter = row*column
        if 0 <= index % iter and index % iter < column:
            vertical_position = vertical
        elif column * (row - 1) <= index % iter and index % iter < iter:
            vertical_position = vertical + margin_height * 2
        else:
            vertical_position = vertical + margin_height
        return vertical_position


    def draw_rectangle(self, image_slide):
        tx_width = 4
        tx_height = 1
        tx_top = Inches((self.ppt_height - tx_height) / 2)
        tx_left = Inches((self.ppt_width - tx_width) / 2)
        rect0 = image_slide.shapes.add_shape(  # shapeオブジェクト➀を追加
            MSO_SHAPE.ROUNDED_RECTANGLE,  # 図形の種類を[丸角四角形]に指定
            tx_left, tx_top,  # 挿入位置の指定　左上の座標の指定
            Inches(tx_width), Inches(tx_height))  # 挿入図形の幅と高さの指定

        rect0.fill.solid()  # shapeオブジェクト➀を単色で塗り潰す
        rect0.fill.fore_color.rgb = RGBColor(250, 100, 100)  # RGB指定で色を指定

        pg = rect0.text_frame.paragraphs[0]  # shapeオブジェクト➀のTextFrameの取得
        pg.text = 'ROUNDED_RECTANGLE'  # TextFrameにテキストを設定
        pg.font.size = Pt(10)  # テキストの文字サイズを10ポイントとする

    # get ratio for both and use the smaller one to ensure that the image would fit in the slide panel
    def get_resize_ratio(self, img_width, img_height, pixel_width, pixel_height):
        ratio_width = pixel_width / img_width
        ratio_height = pixel_height / img_height
        return min(ratio_width, ratio_height)
        return ratio



if __name__ == "__main__":
    main()

