import os, os.path
from pptx import Presentation
from pptx.util import Inches, Pt
import tkinter
from tkinter import filedialog
from tkinter import *
from tkinter import ttk
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
from math import ceil
import sys
import json
import math
import time
import Model
import View
import io


class PPTVariables():
    def __init__(self):
        self.column = 0
        self.row = 0
        self.iter = 0
        self.dpi = 72
        self.width = 0
        self.height = 0
        self.width_in_pixel = 0
        self.height_in_pixel = 0
        self.emus_per_px = 0
        self.panel_pixel_width = 0
        self.panel_pixel_height = 0

    def get_iter(self):
        self.iter = self.column*self.row

    def get_length_in_pixels(self):
        self.width_in_pixel = self.width*self.dpi
        self.height_in_pixel = self.height*self.dpi

    def get_emus_per_px(self):
        self.emus_per_px = int(914400 / self.dpi)

    def get_panel_length(self):
        self.panel_pixel_width = int(self.width_in_pixel / self.column)
        self.panel_pixel_height = int(self.height_in_pixel / self.row)


class Controller():
    def __init__(self):
        self.root = Tk()
        self.model = Model.Model()
        self.view = View.View(self.root)
        self.ppt_variables = PPTVariables()

        self.config = self.load_config()

        self.config_data_into_view()
        self.bind_to_view()
        self.ppt_component = SlideComponents()

    def config_data_into_view(self):
        self.view.input_path_label.configure(text=self.config.input_path)
        self.view.output_path_label.configure(text=self.config.output_path)
        self.view.gui_ppt_name_textbox.delete('0', END)
        self.view.gui_ppt_name_textbox.insert(tkinter.END,self.config.ppt_name)
        self.view.gui_row.delete('0', END)
        self.view.gui_row.insert(tkinter.END,self.config.row)
        self.view.gui_column.delete('0', END)
        self.view.gui_column.insert(tkinter.END, self.config.column)
        self.view.gui_ppt_width.delete('0', END)
        self.view.gui_ppt_width.insert(tkinter.END, self.config.width)
        self.view.gui_ppt_height.delete('0', END)
        self.view.gui_ppt_height.insert(tkinter.END, self.config.height)
        self.view.gui_slide_counter.delete('0', END)
        self.view.gui_slide_counter.insert(tkinter.END, self.config.cell_number)
        self.view.combobox.set(self.config.sort_method)

    def load_config(self):
        config = ConfigObject()
        if os.path.isfile("config.json"):
            with open('config.json', 'r') as f:
                data = json.load(f)
            config = LoadingConfig(data)
        return config

    def bind_to_view(self):
        self.view.input_path_button.bind("<ButtonPress>",
                                         lambda event: self.get_path(event, self.view.input_path_label))
        self.view.output_path_button.bind("<ButtonPress>",
                                          lambda event: self.get_path(event, self.view.output_path_label))
        self.view.start_process_button.bind("<ButtonPress>", lambda event: self.ppt_generation_process(event))
        self.view.save_config_button.bind("<ButtonPress>",
                                          lambda event: self.save_config_into_file(event, self.config))

    def save_config_into_file(self,event,config):
        config.input_path = self.view.input_path_label.cget("text")
        config.output_path = self.view.output_path_label.cget("text")
        config.ppt_name = self.view.gui_ppt_name_textbox.get()
        config.row = self.view.gui_row.get()
        config.column = self.view.gui_column.get()
        config.width = self.view.gui_ppt_width.get()
        config.height = self.view.gui_ppt_height.get()
        config.cell_number = self.view.gui_slide_counter.get()
        config.sort_method = self.view.combobox.get()
        with open('config.json', 'w', encoding='utf-8') as f:
            json.dump(config, f,default=lambda x: x.__dict__, ensure_ascii=False, indent=4)


    def get_path(self, event, config):
        selected_path = filedialog.askdirectory()
        if not selected_path=="":
            config.configure(text=selected_path)

    def ppt_generation_process(self, event):
        #tkinter.Tk().withdraw()
        #input
        self.combobox_value = self.view.combobox.get()
        self.input_path = self.view.input_path_label.cget("text")
        self.img_list = self.get_images(self.input_path)
        self.img_count = len(self.img_list)

        #divide ppt
        self.ppt_variables.row = int(self.view.gui_row.get())
        self.ppt_variables.column = int(self.view.gui_column.get())
        self.ppt_variables.get_iter()

        #ppt dimension
        self.ppt_variables.width = float(self.view.gui_ppt_width.get())
        self.ppt_variables.height = float(self.view.gui_ppt_height.get())
        self.ppt_variables.get_length_in_pixels()

        self.slide_counter = int(self.view.gui_slide_counter.get()) / self.ppt_variables.iter

        self.emus_per_px = self.ppt_variables.get_emus_per_px()

        self.ppt_variables.get_panel_length()
        prs = self.append_images_in_ppt()
        #output
        self.ppt_name = self.view.gui_ppt_name_textbox.get()
        self.output_path = self.view.input_path_label.cget("text")
        prs.save(self.output_path + '/' + self.ppt_name + '.pptx')
        os.startfile(self.output_path + '/' + self.ppt_name + '.pptx')


    def run(self):
        self.root.minsize(width=500, height=300)
        self.root.title("Image2ppt")
        self.root.mainloop()

    def sort_images_alphabetically(self,list_of_files,reverse=False):
        return sorted(list_of_files,reverse=reverse)

    def sort_images_by_date(self,list_of_files,dir_name,reverse=False):
        list_of_files = sorted(list_of_files,
                               key=lambda x: os.path.getmtime(os.path.join(dir_name, x))
                               , reverse=reverse
                               )
        return list_of_files

    def get_list_of_files(self,dir_name):
        list_of_files = filter(lambda x: os.path.isfile(os.path.join(dir_name, x)),
               os.listdir(dir_name))
        return list_of_files

    def get_images(self, input_path):
        # folder_files = os.listdir(input_path)
        img_list = []
        list_of_files = self.get_list_of_files(input_path)

        if self.combobox_value == 'Alphabetical A-Z':
            list_of_files = self.sort_images_alphabetically(list_of_files)
        elif self.combobox_value == 'Alphabetical Z-A':
            list_of_files = self.sort_images_alphabetically(list_of_files,True)
        elif self.combobox_value == "Oldest-Newest":
            list_of_files = self.sort_images_by_date(list_of_files, input_path)
        elif self.combobox_value == "Newest-Oldest":
            list_of_files = self.sort_images_by_date(list_of_files, input_path,True)

        for file_name in list_of_files:
            self.model.check_image_extension(img_list,file_name)

        return img_list



    # generate ppt and add images to the ppt
    def append_images_in_ppt(self):
        prs = Presentation()
        prs.slide_width = Inches(self.ppt_variables.width)
        prs.slide_height = Inches(self.ppt_variables.height)
        prs = self.append_images(prs)
        return prs

    def append_images(self, prs):
        blank_slide = prs.slide_layouts[6]
        #start adding images on the slide

        for i in range(self.img_count):
            #prepare blank slide if the image reaches threshold
            if i % self.ppt_variables.iter == 0:
                image_slide = prs.slides.add_slide(blank_slide)
                #cell_number_textbox_visible= False;
                #if cell_number_textbox_visible:
                #    cell_number = str(ceil(len(prs.slides)/self.slide_counter))
                #    self.ppt_component.textbox(image_slide,cell_number,self.ppt_variables.width,self.ppt_variables.height)


            #prepare image
            current_img = Image.open(self.input_path + '/' + self.img_list[i])

            #resize image
            ratio = self.model.get_resize_ratio(current_img.width, current_img.height, self.ppt_variables.panel_pixel_width, self.ppt_variables.panel_pixel_height)
            resized_img = current_img.resize((int(current_img.width * ratio), int(current_img.height * ratio)))

            margin_width = self.model.get_margin_in_pixel(self.ppt_variables.panel_pixel_width, resized_img.width)
            margin_height = self.model.get_margin_in_pixel(self.ppt_variables.panel_pixel_height, resized_img.height)

            #prepare panel location
            prefixed_horizontal_location = (i % self.ppt_variables.column) * (self.ppt_variables.width_in_pixel / self.ppt_variables.column)
            prefixed_vertical_location = (i % self.ppt_variables.iter // self.ppt_variables.column) *self.ppt_variables.height_in_pixel / self.ppt_variables.row

            #prepare image location based on panel location
            fixed_vertical_position = self.model.apply_vertical_margin(i, self.ppt_variables.row, self.ppt_variables.column, prefixed_vertical_location, margin_height)
            fixed_horizontal_position = prefixed_horizontal_location + margin_width

            #add image to a panel
            with io.BytesIO() as output:
                #resized_img.save(output, format="GIF")
                quality_val = 100
                resized_img.save(output,format = "GIF",quality=quality_val)
                image_slide.shapes.add_picture(output, fixed_horizontal_position*self.ppt_variables.emus_per_px, fixed_vertical_position*self.ppt_variables.emus_per_px)

        self.ppt_component.draw_rectangle(image_slide,self.ppt_variables.width,self.ppt_variables.height)

        return prs


class LoadingConfig(object):
    def __init__(self, dict):
        vars(self).update(dict)

class ConfigObject:
    def __init__(self):
        self.input_path = "Please choose an input folder"
        self.output_path = "Please choose an output folder"
        self.ppt_name = "test"

        self.row = 2
        self.column = 4
        self.width = 26.6666
        self.height = 15
        self.cell_number = 16
        self.sort_method = 'Alphabetical A-Z'

class SlideComponents:
    def textbox(self, image_slide,cell_number,ppt_width,ppt_height):

        central_box = image_slide.shapes.add_textbox(Inches(ppt_width / 2 - 0.65),
                                                     Inches(ppt_height / 2 - 0.6), Inches(1), Inches(1))
        central_label = central_box.text_frame.add_paragraph()
        central_label.text = "Cell " + cell_number
        central_label.font.size = Pt(30)


    def draw_rectangle(self, image_slide,ppt_width,ppt_height):
        tx_width = 4
        tx_height = 1
        tx_top = Inches((ppt_height - tx_height) / 2)
        tx_left = Inches((ppt_width - tx_width) / 2)
        rect0 = image_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            tx_left, tx_top,
            Inches(tx_width), Inches(tx_height))

        rect0.fill.solid()
        rect0.fill.fore_color.rgb = RGBColor(250, 100, 100)

        pg = rect0.text_frame.paragraphs[0]
        pg.text = 'ROUNDED_RECTANGLE'
        pg.font.size = Pt(10)

