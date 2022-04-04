# import unittest
# import Image2ppt
#
# class MyTestCase(unittest.TestCase):
#     def test_addition(self):
#         add = Image2ppt.AddTest()
#         result = add.add_test(1,2)
#         self.assertEqual(3, result)  # add assertion here
#
#
# if __name__ == '__main__':
#     unittest.main()


from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

#tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

prs.save('test.pptx')
