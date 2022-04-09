import os, os.path
from pptx import Presentation
from pptx.util import Inches, Pt
import tkinter
from tkinter import filedialog
from tkinter import *
from tkinter import ttk
import io
from PIL import Image
from math import ceil
import math

# Utilizes python-pptx: https://python-pptx.readthedocs.io/


def main():

    create_gui = CreateGUI()
    create_gui.construct_form()


class Components:

    def create_textbox(self,frame,text,row,column):
        txt = ttk.Entry(frame,width=40)
        txt.insert(tkinter.END,text)
        txt.grid(row=row,column=column)
        return txt

    def create_label(self,frame,text,row,column):
        label1 = ttk.Label(
            frame,
            text=text,
            padding=(5, 10))
        label1.grid(row=row, column=column)
        return label1

    def create_button(self,frame, text, row, column):
        button1 = ttk.Button(
            frame,
            text=text,
            )
        button1.grid(row=row, column=column)
        return button1

    def create_notebook(self, root):
        notebook = ttk.Notebook(root)
        notebook.grid()
        return notebook

    def create_gentab(self, notebook):
        frame1 = ttk.Frame(notebook)
        notebook.add(frame1, text="General")
        return frame1

    def create_advtab(self, notebook):
        frame2 = ttk.Frame(notebook)
        notebook.add(frame2, text="Advanced")
        return frame2
    # we can merge tab creation into one by adding text input
    # if we put "self" in front of a variable then, the variable becomes global within the class
    #but if the variable does not have self then it becomes a local variable
    # in this case, we only need to have the variable "frame" within the function, so we dont need to put self

    def create_tab(self,notebook,text):
        frame = ttk.Frame(notebook)
        notebook.add(frame, text=text)
        return frame

class CreateGUI:
    def __init__(self):

        self.input_path_label = "initial path"
        self.output_path_label = "initial path"
        self.ppt_name_textbox = "test"
        self.path_list = [r"C:\Users\Spencer Laptop\Documents\py\Image2ppt__\Image2ppt\Input",
                          r"C:\Users\Spencer Laptop\Documents\py\Image2ppt__\Image2ppt\Output"]

    def construct_form(self):
        components = Components()
        root = Tk()
        root.minsize(width=500, height=300)
        root.title("Image2ppt")
        notebook = components.create_notebook(root)

        #general and advanced tabs come from the same function
        # also renamed frame1 and frame2
        #if you right click on a variable and go to the refactor then you can rename variable at once.
        frame_general = components.create_tab(notebook,"General")

        self.input_path_label = components.create_label(frame_general, self.path_list[0], 0, 0)

        input_path_button = components.create_button(frame_general, "Input path", 0, 1)
        input_path_button.bind("<ButtonPress>", lambda event: self.get_path(event, self.input_path_label))

        self.output_path_label = components.create_label(frame_general, self.path_list[1], 1, 0)

        output_path_button = components.create_button(frame_general,"Output path", 1, 1)
        output_path_button.bind("<ButtonPress>", lambda event: self.get_path(event, self.output_path_label))

        gui_column_desc = components.create_label(frame_general, "Column Number:", 2, 0)
        self.gui_column = components.create_textbox(frame_general, 4, 2, 1)

        gui_row_desc = components.create_label(frame_general, "Row Number:", 3, 0)
        self.gui_row = components.create_textbox(frame_general, 2, 3, 1)

        frame_advanced = components.create_tab(notebook, "Advanced")

        gui_ppt_width_desc = components.create_label(frame_advanced, "Slide Width:", 4, 0)
        self.gui_ppt_width = components.create_textbox(frame_advanced, 13.333, 4, 1)

        gui_ppt_height_desc = components.create_label(frame_advanced, "Slide Height:", 5, 0)
        self.gui_ppt_height = components.create_textbox(frame_advanced, 7.5, 5, 1)

        gui_slide_counter_desc = components.create_label(frame_advanced, "Images for each cell:", 6, 0)
        self.gui_slide_counter = components.create_textbox(frame_advanced, 16, 6, 1)

        ppt_name_label = components.create_label (frame_general, "Save Name:", 7, 0)
        self.gui_ppt_name_textbox = components.create_textbox(frame_general, "test", 7, 1)

        start_process_button = components.create_button(frame_general, "Start", 8, 1)
        start_process_button.bind("<ButtonPress>", lambda event: self.ppt_generation_process(event))

        self.start_gui(root)


    def get_path(self, event, arg):
        selected_path = filedialog.askdirectory()
        if not selected_path=="":
            arg.configure(text=selected_path)

    def ppt_generation_process(self, event):
        #tkinter.Tk().withdraw()
        path_list = [self.input_path_label.cget("text"), self.output_path_label.cget("text")]
        input_path = path_list[0]
        output_path = path_list[1]
        parameters = [self.gui_column.get(), self.gui_row.get(), self.gui_slide_counter.get(), self.gui_ppt_width.get(), self.gui_ppt_height.get()]

        append_slide = AppendSlide(input_path)
        append_slide.get_parameters(parameters)
        prs = append_slide.append_images_in_ppt()
        prs.save(output_path + '/' + self.gui_ppt_name_textbox.get() + '.pptx')
        os.startfile(output_path + '/' + self.gui_ppt_name_textbox.get() + '.pptx')

    def start_gui(self,root):
        root.mainloop()


#this is not used in production
class AddTest:
    def add_test(self,a,b):
        return a+b


class AppendSlide:
    # initial value loading
    def __init__(self, input_path):
        self.img_list = self.get_images(input_path)
        self.img_count = len(self.img_list)

        self.input_path = input_path
        self.slide_number = 1

    def get_parameters(self, parameters):
        self.column = int(parameters[0])
        self.row =  int(parameters[1])
        self.ppt_width =  float(parameters[3])
        self.ppt_height = float(parameters[4])
        self.img_width = self.ppt_width / self.column
        self.img_height = self.ppt_height / self.row
        self.img_iter = self.column * self.row
        self.slide_counter = int(parameters[2]) / self.img_iter

    def get_images(self, input_path):
        folder_files = os.listdir(input_path)
        img_list = []

        for file in folder_files:
            if file.endswith('.png') or file.endswith(".tif") or file.endswith(".jpg") or file.endswith(".jpeg"):
                img_list.append(file)

        print(img_list)
        return img_list

    # generate ppt and add images to the ppt
    def append_images_in_ppt(self):
        prs = Presentation()
        prs.slide_width = Inches(self.ppt_width)
        prs.slide_height = Inches(self.ppt_height)
        prs = self.append_images(prs)
        return prs

    def append_images(self, prs):
        blank_slide = prs.slide_layouts[6]
        #in pixels
        pixel_width = int(960 / self.column)
        pixel_height = int(540 / self.row)
        #in inches
        width = self.ppt_width / self.column
        height = self.ppt_height / self.row
        dpi = 72

        for i in range(self.img_count):
            if i % self.img_iter == 0:
                image_slide = prs.slides.add_slide(blank_slide)
                cellnumber = str(ceil(self.slide_number/self.slide_counter))
                central_box = image_slide.shapes.add_textbox(Inches(self.ppt_width/2-0.65), Inches(self.ppt_height/2-0.6), Inches(1), Inches(1))
                central_label = central_box.text_frame.add_paragraph()
                central_label.text = "Cell " + cellnumber
                central_label.font.size = Pt(30)
                self.slide_number += 1
                print(self.slide_number)

            #in inches
            horizontal = Inches((i % self.column) * (self.ppt_width / self.column))
            vertical = Inches((i % self.img_iter // self.column) * self.ppt_height / self.row)

            current_img = Image.open(self.input_path + '/' + self.img_list[i])
            ratio = self.get_resize_ratio(current_img,pixel_width,pixel_height)
            resized_img = current_img.resize((int(current_img.width * ratio), int(current_img.height * ratio)))

            margin_width = Inches(width-resized_img.width/dpi)/2
            margin_height = Inches(height-resized_img.height/dpi)/2
            if i % self.img_iter // self.column == 0:
                vertical = Inches(0)
            elif i % self.img_iter // self.column == self.row -1:
                vertical = Inches(self.ppt_height - resized_img.height/72)
            else:
                vertical = vertical + margin_height

            with io.BytesIO() as output:
                resized_img.save(output, format="GIF")
                image_slide.shapes.add_picture(output, horizontal+margin_width, vertical)


        return prs

    #resize ratio is determined depending on the orientation of the image
    #if we have a perfect square,we make sure that the picture fits in the slide by choosing the smaller side of the slide compartment
    def get_resize_ratio(self,img,pixel_width,pixel_height):
        if img.width > img.height:
            ratio = pixel_width / img.width
        elif img.width < img.height:
            ratio = pixel_height / img.height
        elif img.width == img.height:
            ratio = min(pixel_width, pixel_height) / img.width
        return ratio



if __name__ == "__main__":
    main()
